from flask import Flask, render_template, request, jsonify
import os
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
from dotenv import load_dotenv
from openai import AzureOpenAI
import requests
from pathlib import Path
import PyPDF2
from pptx import Presentation
from docx import Document
import logging
import nltk
from nltk.tokenize import sent_tokenize
import re
from rank_bm25 import BM25Okapi
from sentence_transformers.cross_encoder import CrossEncoder

nltk.download('punkt', quiet=True)
nltk.download('punkt_tab', quiet=True)

# 配置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

load_dotenv()

app = Flask(__name__)

# API 配置
API_TYPE = os.environ.get("API_TYPE")  
DEEPSEEK_API_KEY = os.environ.get("DEEPSEEK_API_KEY")
DEEPSEEK_ENDPOINT = os.environ.get("DEEPSEEK_ENDPOINT")
AZURE_OPENAI_KEY = os.environ.get("AZURE_OPENAI_KEY")
AZURE_OPENAI_ENDPOINT = os.environ.get("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_VERSION = os.environ.get("AZURE_OPENAI_VERSION")
AZURE_OPENAI_DEPLOYMENT = os.environ.get("AZURE_OPENAI_DEPLOYMENT")

# 验证 API 配置
if API_TYPE == "deepseek" and (not DEEPSEEK_API_KEY or not DEEPSEEK_ENDPOINT):
    logger.error("DeepSeek API configuration missing")
    raise ValueError("DEEPSEEK_API_KEY and DEEPSEEK_ENDPOINT must be set")
elif API_TYPE == "azure" and not all([AZURE_OPENAI_KEY, AZURE_OPENAI_ENDPOINT, AZURE_OPENAI_VERSION, AZURE_OPENAI_DEPLOYMENT]):
    logger.error("Azure OpenAI configuration missing")
    raise ValueError("All Azure OpenAI configurations must be set")
elif API_TYPE not in ["deepseek", "azure"]:
    logger.error(f"Invalid API_TYPE: {API_TYPE}")
    raise ValueError("API_TYPE must be 'deepseek' or 'azure'")

if API_TYPE == "azure":
    azure_client = AzureOpenAI(
        api_key=AZURE_OPENAI_KEY,
        api_version=AZURE_OPENAI_VERSION,
        azure_endpoint=AZURE_OPENAI_ENDPOINT
    )

# 初始化模型
embedder = SentenceTransformer('intfloat/e5-large')
cross_encoder = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-12-v2')

# 全局变量
documents = []
embeddings = []
chunk_metadata = []
faiss_index = None
bm25_index = None
INDEX_INITIALIZED = False
conversation_history = []

def detect_clauses(text):
    clause_pattern = r'\b(?:Section|Clause)?\s*(\d+(?:[.-]\d+)*\s*(?:[A-Z]|\([a-zA-Z]\))?|[A-Z](?:[.-]\d+)*\s*(?:[a-z]|\([a-z]\))?)\s*[-:]?\s*'
    matches = list(re.finditer(clause_pattern, text))
    clauses = []
    for i, match in enumerate(matches):
        clause_id = match.group(1).strip()
        start_pos = match.start()
        next_start = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        content = text[start_pos:next_start].strip()
        if content:
            clauses.append({'clause_id': clause_id, 'content': content, 'start_pos': start_pos})
    logger.debug(f"Detected {len(clauses)} clauses")
    return clauses

def chunk_text(text, target_size=3000, min_size=1500, overlap=500):
    if not text.strip():
        return []
    
    clauses = detect_clauses(text)
    chunks = []
    start_pos = 0
    text_length = len(text)
    
    while start_pos < text_length:
        end_pos = min(start_pos + target_size, text_length)
        chunk_text = text[start_pos:end_pos]
        
        last_clause_end = start_pos
        for clause in clauses:
            if clause['start_pos'] >= start_pos and clause['start_pos'] < end_pos:
                last_clause_end = max(last_clause_end, clause['start_pos'] + len(clause['content']))
            elif clause['start_pos'] >= end_pos:
                break
        
        if last_clause_end > start_pos and last_clause_end - start_pos >= min_size:
            end_pos = last_clause_end
        elif end_pos - start_pos < min_size and end_pos < text_length:
            for clause in clauses:
                if clause['start_pos'] > end_pos:
                    end_pos = clause['start_pos'] + len(clause['content'])
                    break
        
        chunk = text[start_pos:end_pos].strip()
        if chunk:
            chunk_clauses = [c for c in clauses if c['start_pos'] >= start_pos and c['start_pos'] < end_pos]
            chunks.append(chunk)
            logger.debug(f"Chunk created: {len(chunk)} chars, clauses: {len(chunk_clauses)}")
        
        start_pos = end_pos - overlap if end_pos < text_length else end_pos
    
    logger.info(f"Chunked text into {len(chunks)} chunks")
    return chunks

def extract_text_from_docx(file_path):
    logger.info(f"Reading DOCX: {file_path}")
    try:
        doc = Document(file_path)
        text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
        if text:
            chunks = chunk_text(text)
            clauses = detect_clauses(text)
            return [{
                'text': chunk,
                'page': None,
                'clause_ids': [c['clause_id'] for c in clauses if c['clause_id'] in chunk]
            } for chunk in chunks]
        return []
    except Exception as e:
        logger.error(f"Error reading DOCX {file_path}: {str(e)}")
        return []

def extract_text_from_pdf(file_path):
    logger.info(f"Reading PDF: {file_path}")
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            chunks_with_metadata = []
            for page_num, page in enumerate(pdf_reader.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    chunks = chunk_text(text)
                    clauses = detect_clauses(text)
                    chunks_with_metadata.extend({
                        'text': chunk,
                        'page': page_num,
                        'clause_ids': [c['clause_id'] for c in clauses if c['clause_id'] in chunk]
                    } for chunk in chunks)
            return chunks_with_metadata
    except Exception as e:
        logger.error(f"Error reading PDF {file_path}: {str(e)}")
        return []

def extract_text_from_ppt(file_path):
    logger.info(f"Reading PPT: {file_path}")
    try:
        prs = Presentation(file_path)
        chunks_with_metadata = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
            if text:
                chunks = chunk_text(text)
                clauses = detect_clauses(text)
                chunks_with_metadata.extend({
                    'text': chunk,
                    'page': slide_num,
                    'clause_ids': [c['clause_id'] for c in clauses if c['clause_id'] in chunk]
                } for chunk in chunks)
        return chunks_with_metadata
    except Exception as e:
        logger.error(f"Error reading PPT {file_path}: {str(e)}")
        return []

def extract_text_from_txt(file_path):
    logger.info(f"Reading TXT: {file_path}")
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
            if text.strip():
                chunks = chunk_text(text)
                clauses = detect_clauses(text)
                return [{
                    'text': chunk,
                    'page': None,
                    'clause_ids': [c['clause_id'] for c in clauses if c['clause_id'] in chunk]
                } for chunk in chunks]
        return []
    except Exception as e:
        logger.error(f"Error reading TXT {file_path}: {str(e)}")
        return []

def initialize_index():
    global documents, embeddings, chunk_metadata, faiss_index, bm25_index, INDEX_INITIALIZED
    if INDEX_INITIALIZED:
        logger.info("Index already initialized")
        return
    
    documents_dir = "documents"
    if not os.path.exists(documents_dir):
        os.makedirs(documents_dir)
        logger.warning(f"Created empty directory: {documents_dir}")
    
    documents.clear()
    embeddings.clear()
    chunk_metadata.clear()
    
    file_extractors = {
        '.txt': extract_text_from_txt,
        '.pdf': extract_text_from_pdf,
        '.pptx': extract_text_from_ppt,
        '.docx': extract_text_from_docx
    }
    
    logger.info("Starting document indexing...")
    for root, _, files in os.walk(documents_dir):
        for filename in files:
            ext = Path(filename).suffix.lower()
            if ext in file_extractors:
                filepath = os.path.join(root, filename)
                chunks_with_metadata = file_extractors[ext](filepath)
                for i, chunk_info in enumerate(chunks_with_metadata):
                    documents.append(chunk_info['text'])
                    chunk_metadata.append({
                        'filepath': filepath,
                        'chunk_index': i,
                        'text': chunk_info['text'],
                        'page': chunk_info['page'],
                        'clause_ids': chunk_info['clause_ids']
                    })
    
    if documents:
        logger.info(f"Generating embeddings for {len(documents)} chunks")
        embeddings = embedder.encode(documents, show_progress_bar=True, normalize_embeddings=True)
        embeddings = np.array(embeddings).astype('float32')
        dimension = embeddings.shape[1]
        quantizer = faiss.IndexFlatL2(dimension)
        faiss_index = faiss.IndexIVFFlat(quantizer, dimension, 50)
        faiss_index.train(embeddings)
        faiss_index.add(embeddings)
        faiss_index.nprobe = 10
        
        tokenized_corpus = [doc.lower().split() for doc in documents]
        bm25_index = BM25Okapi(tokenized_corpus)
        logger.info(f"Indexed {faiss_index.ntotal} vectors and BM25 corpus")
    else:
        logger.warning("No documents indexed")
        faiss_index = None
        bm25_index = None
    
    INDEX_INITIALIZED = True

initialize_index()

prompt_template = """Answer based solely on the following documents. Do not use external knowledge. Structure your response:
1. Summary paragraph.
2. Detailed information (use bullet points only for lists).
3. Cite Reference IDs ([RefX]) inline for document info used; state "No documents used" if none apply.
4. List references at the end (e.g., References: [Ref1]).

Documents:
{passages}

Query: {query}

Answer:"""

def create_prompt(passages, query):
    passages_text = "\n".join(f"[Ref{idx+1}] {passage}" for idx, passage in enumerate(passages)) if passages else "No relevant documents found."
    prompt = prompt_template.format(passages=passages_text, query=query)
    logger.debug(f"Prompt passages: {passages_text[:100]}...")
    return prompt, [(f"[Ref{idx+1}]", passage) for idx, passage in enumerate(passages)]

def format_response(bot_message, passage_refs, chunks):
    lines = bot_message.strip().split('\n')
    formatted = []
    in_list = False
    used_refs = set(re.findall(r'\[Ref\d+\]', bot_message))
    
    for line in lines:
        line = line.strip()
        if line.startswith("References:"):
            continue
        line = re.sub(r'\[Ref\d+\]', '', line).strip()
        if not line:
            if in_list:
                formatted.append('</ul>')
                in_list = False
            continue
        if line.startswith('- '):
            if not in_list:
                formatted.append('<ul>')
                in_list = True
            formatted.append(f'<li>{line[2:]}</li>')
        else:
            if in_list:
                formatted.append('</ul>')
                in_list = False
            formatted.append(f'<p>{line}</p>')
    
    if in_list:
        formatted.append('</ul>')
    
    citations = []
    for ref_id, passage in passage_refs:
        if ref_id in used_refs:
            filepath = passage.split('\n')[0].replace("Document: ", "")
            citation = f"Source: {os.path.basename(filepath)}"
            pages = next((c['page'] for c in chunks if c['filepath'] == filepath and c['page']), None)
            if pages:
                citation += f", Page: {pages}"
            citations.append(citation)
    
    if citations:
        formatted.append(f'<div><strong>References:</strong><br>{"<br>".join(citations)}</div>')
    
    return ''.join(formatted)

def parse_query_clause(query):
    clause_pattern = r'\b(?:Section|Clause)?\s*(\d+(?:[.-]\d+)*\s*(?:[A-Z]|\([a-zA-Z]\))?|[A-Z](?:[.-]\d+)*\s*(?:[a-z]|\([a-z]\))?)\b'
    match = re.search(clause_pattern, query)
    return match.group(1).strip() if match else None

def find_related_clauses(clause_id, clause_ids_list):
    related = []
    for clause_ids in clause_ids_list:
        for cid in clause_ids:
            if clause_id in cid or cid.startswith(clause_id + '.') or cid.startswith(clause_id + '-'):
                related.append(cid)
    return related

def enhance_query(query, clause_id):
    if clause_id:
        return f"What is the content of clause {clause_id} in the document?"
    return query

@app.route('/')
def home():
    global conversation_history
    conversation_history = []
    logger.info("Conversation history reset")
    return render_template('chat.html')

@app.route('/chat', methods=['POST'])
def chat():
    global conversation_history
    
    # 获取用户消息
    user_message = request.json.get('message', '').strip()
    logger.info(f"Received request: {request.json}")
    if not user_message:
        logger.warning("No user message provided")
        response = "<p>Error: No message provided. Please enter a query.</p>"
        return jsonify({"message": response})
    
    conversation_history.append({"role": "user", "content": user_message})
    logger.info(f"User message added: {user_message}, History length: {len(conversation_history)}")
    
    # 初始化默认值
    passages = ["No relevant content found in documents."]
    chunks = []
    
    # 检索阶段
    try:
        query_clause = parse_query_clause(user_message)
        enhanced_query = enhance_query(user_message, query_clause)
        logger.info(f"Query: {user_message}, Clause: {query_clause}, Enhanced: {enhanced_query}")
        
        query_embedding = embedder.encode([enhanced_query], show_progress_bar=False, normalize_embeddings=True).astype('float32')
        query_embedding = query_embedding.reshape(1, -1)
        
        if faiss_index and bm25_index and faiss_index.ntotal > 0:
            distances, indices = faiss_index.search(query_embedding, k=20)
            faiss_scores = 1 - (distances[0] / (np.max(distances[0]) + 1e-6))
            faiss_results = {i: score for i, score in zip(indices[0], faiss_scores) if 0 <= i < len(chunk_metadata)}
            logger.info(f"FAISS results: {len(faiss_results)} items, top scores: {list(faiss_results.values())[:5]}")
            
            tokenized_query = user_message.lower().split()
            bm25_scores = bm25_index.get_scores(tokenized_query)
            bm25_max = max(bm25_scores) + 1e-6
            bm25_results = {i: score / bm25_max for i, score in enumerate(bm25_scores) if score > 0}
            logger.info(f"BM25 results: {len(bm25_results)} items, top scores: {list(bm25_results.values())[:5]}")
            
            hybrid_scores = {}
            faiss_weight = 0.3 if query_clause else 0.7
            bm25_weight = 0.7 if query_clause else 0.3
            for i in set(faiss_results.keys()).union(bm25_results.keys()):
                faiss_score = faiss_results.get(i, 0)
                bm25_score = bm25_results.get(i, 0)
                hybrid_scores[i] = faiss_weight * faiss_score + bm25_weight * bm25_score
            
            sorted_indices = sorted(hybrid_scores.keys(), key=lambda i: hybrid_scores[i], reverse=True)[:20]
            logger.info(f"Hybrid top indices: {sorted_indices[:5]}, scores: {[hybrid_scores[i] for i in sorted_indices[:5]]}")
            
            rerank_pairs = [(enhanced_query, chunk_metadata[i]['text']) for i in sorted_indices]
            rerank_scores = cross_encoder.predict(rerank_pairs)
            reranked_indices = [sorted_indices[i] for i in np.argsort(rerank_scores)[::-1]]
            logger.info(f"Reranked indices: {reranked_indices[:5]}, scores: {rerank_scores[np.argsort(rerank_scores)[::-1]][:5]}")
            
            selected_indices = set()
            if query_clause:
                related_clauses = find_related_clauses(query_clause, [m['clause_ids'] for m in chunk_metadata])
                logger.info(f"Related clauses: {related_clauses}")
                exact_matches = [i for i in reranked_indices if query_clause in chunk_metadata[i]['clause_ids']]
                related_matches = [i for i in reranked_indices if any(cid in chunk_metadata[i]['clause_ids'] for cid in related_clauses)]
                
                if exact_matches:
                    for i in exact_matches[:3]:
                        selected_indices.add(i)
                        if i > 0:
                            selected_indices.add(i - 1)
                        if i + 1 < len(chunk_metadata):
                            selected_indices.add(i + 1)
                elif related_matches:
                    for i in related_matches[:3]:
                        selected_indices.add(i)
                        if i > 0:
                            selected_indices.add(i - 1)
                        if i + 1 < len(chunk_metadata):
                            selected_indices.add(i + 1)
                else:
                    selected_indices.update(reranked_indices[:5])
            else:
                selected_indices.update(reranked_indices[:7])
            
            logger.info(f"Selected indices: {list(selected_indices)}")
            passages = []
            chunks = []
            for i in sorted(selected_indices):
                chunk = chunk_metadata[i]
                chunks.append(chunk)
                clause_info = f"Clauses: {', '.join(chunk['clause_ids'])}" if chunk['clause_ids'] else "No clauses"
                passage = f"Document: {chunk['filepath']}\n{clause_info}\nContent: {chunk['text']}"
                passages.append(passage)
        else:
            logger.warning("Index not initialized or empty")
    except Exception as e:
        logger.error(f"Retrieval error: {str(e)}")
        passages = ["Retrieval error occurred."]
    
    # 生成 Prompt
    prompt, passage_refs = create_prompt(passages, user_message)
    logger.info(f"Prompt generated: {prompt[:100]}...")
    
    # 构造 Messages
    system_message = {"role": "system", "content": "You are an assistant to the project manager, specializing in retrieving and analyzing information from internal project documents. Your role includes identifying risks (with severity: Critical/High/Medium/Low), flagging client clarification needs, extracting requirements and dependencies (including timelines), and detailing concrete quantitative data (e.g., numbers, volumes, sizes) in scope. When requested, provide estimations using document data first, supplemented by pre-trained knowledge if needed, and label assumptions clearly. Provide accurate answers based solely on documents unless estimations require otherwise, avoiding speculation. Structure responses with a brief summary followed by detailed analysis (e.g., Risks, Quantities, Dependencies, Clarifications, Estimations if requested). Cite Reference IDs ([RefX]) for document info; state if no documents are used. List references at the end (e.g., 'References: [Ref1]') and note 'Used pre-trained knowledge for estimation' if applicable."}
    messages = [system_message] + conversation_history + [{"role": "user", "content": prompt}]
    logger.info(f"Messages constructed: {len(messages)} items, first message: {messages[0]}")
    
    # API 调用
    try:
        if API_TYPE == "deepseek":
            headers = {"Authorization": f"Bearer {DEEPSEEK_API_KEY}", "Content-Type": "application/json"}
            payload = {"model": "deepseek-chat", "messages": messages, "max_tokens": 1500, "temperature": 0.7}
            logger.info(f"Sending DeepSeek request: {payload['messages'][0]}...")
            response = requests.post(f"{DEEPSEEK_ENDPOINT}/chat/completions", headers=headers, json=payload, timeout=30)
            response.raise_for_status()
            bot_message = response.json()["choices"][0]["message"]["content"].strip()
        elif API_TYPE == "azure":
            logger.info(f"Sending Azure request: {messages[0]}...")
            response = azure_client.chat.completions.create(
                model=AZURE_OPENAI_DEPLOYMENT, messages=messages, max_tokens=1500, temperature=0.7
            )
            bot_message = response.choices[0].message.content.strip()
        else:
            logger.error("Invalid API_TYPE")
            bot_message = "API configuration error."
        
        formatted_message = format_response(bot_message, passage_refs, chunks)
        conversation_history.append({"role": "assistant", "content": bot_message})
        logger.info(f"Response received: {bot_message[:100]}...")
    except Exception as e:
        logger.error(f"API call failed: {str(e)}")
        bot_message = f"Error occurred while processing your request: {str(e)}"
        formatted_message = f"<p>{bot_message}</p>"
    
    return jsonify({"message": formatted_message})

if __name__ == '__main__':
    app.run(debug=True)