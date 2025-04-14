from flask import Flask, render_template, request, jsonify
import os
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
from dotenv import load_dotenv
from openai import AzureOpenAI
import requests
from pathlib import Path
import PyPDF2
from pptx import Presentation
from docx import Document
import logging
import nltk
from nltk.tokenize import sent_tokenize
from sklearn.cluster import KMeans
import re
import yaml  

# Download NLTK data for sentence tokenization
nltk.download('punkt', quiet=True)
nltk.download('punkt_tab', quiet=True)

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

app = Flask(__name__)

# Retrieve configurations from environment variables
API_TYPE = os.environ.get("API_TYPE")
DEEPSEEK_API_KEY = os.environ.get("DEEPSEEK_API_KEY")
DEEPSEEK_ENDPOINT = os.environ.get("DEEPSEEK_ENDPOINT")
AZURE_OPENAI_KEY = os.environ.get("AZURE_OPENAI_KEY")
AZURE_OPENAI_ENDPOINT = os.environ.get("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_VERSION = os.environ.get("AZURE_OPENAI_VERSION")
AZURE_OPENAI_DEPLOYMENT = os.environ.get("AZURE_OPENAI_DEPLOYMENT")

# Validate configurations
if API_TYPE == "deepseek":
    if not DEEPSEEK_API_KEY or not DEEPSEEK_ENDPOINT:
        raise ValueError("DEEPSEEK_API_KEY and DEEPSEEK_ENDPOINT must be set")
elif API_TYPE == "azure":
    if not all([AZURE_OPENAI_KEY, AZURE_OPENAI_ENDPOINT, AZURE_OPENAI_VERSION, AZURE_OPENAI_DEPLOYMENT]):
        raise ValueError("All Azure OpenAI configurations must be set")
else:
    raise ValueError("API_TYPE must be either 'deepseek' or 'azure'")

# Initialize Azure OpenAI client if using Azure
if API_TYPE == "azure":
    azure_client = AzureOpenAI(
        api_key=AZURE_OPENAI_KEY,
        api_version=AZURE_OPENAI_VERSION,
        azure_endpoint=AZURE_OPENAI_ENDPOINT
    )

# Initialize embedding model
embedder = SentenceTransformer('multi-qa-mpnet-base-dot-v1')

# Global variables for indexing
documents = []
embeddings = []
chunk_metadata = []
faiss_index = None
INDEX_INITIALIZED = False

# Store conversation history (for multi-turn dialogue)
conversation_history = []

# Document processing functions with page tracking
def extract_text_from_pdf(file_path):
    logger.info(f"Attempting to read PDF: {file_path}")
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            chunks_with_pages = []
            for page_num, page in enumerate(pdf_reader.pages, start=1):
                text = page.extract_text() + "\n"
                if text.strip():
                    chunks = chunk_text(text)
                    for chunk in chunks:
                        chunks_with_pages.append({
                            'text': chunk,
                            'page': page_num
                        })
            logger.info(f"Successfully read {len(chunks_with_pages)} chunks from {file_path}")
            return chunks_with_pages
    except Exception as e:
        logger.error(f"Error reading PDF {file_path}: {str(e)}")
        return []

def extract_text_from_ppt(file_path):
    logger.info(f"Attempting to read PPT: {file_path}")
    try:
        prs = Presentation(file_path)
        chunks_with_pages = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            if text.strip():
                chunks = chunk_text(text)
                for chunk in chunks:
                    chunks_with_pages.append({
                        'text': chunk,
                        'page': slide_num
                    })
        logger.info(f"Successfully read {len(chunks_with_pages)} chunks from {file_path}")
        return chunks_with_pages
    except Exception as e:
        logger.error(f"Error reading PPT {file_path}: {str(e)}")
        return []

def extract_text_from_docx(file_path):
    logger.info(f"Attempting to read DOCX: {file_path}")
    try:
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        if text.strip():
            chunks = chunk_text(text)
            chunks_with_pages = [{'text': chunk, 'page': None} for chunk in chunks]
            logger.info(f"Successfully read {len(chunks_with_pages)} chunks from {file_path}")
            return chunks_with_pages
        return []
    except Exception as e:
        logger.error(f"Error reading DOCX {file_path}: {str(e)}")
        return []

def extract_text_from_txt(file_path):
    logger.info(f"Attempting to read TXT: {file_path}")
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
            if content.strip():
                chunks = chunk_text(content)
                chunks_with_pages = [{'text': chunk, 'page': None} for chunk in chunks]
                logger.info(f"Successfully read {len(chunks_with_pages)} chunks from {file_path}")
                return chunks_with_pages
        return []
    except UnicodeDecodeError:
        logger.warning(f"UTF-8 decoding failed for {file_path}, trying latin-1")
        try:
            with open(file_path, 'r', encoding='latin-1') as file:
                content = file.read()
                if content.strip():
                    chunks = chunk_text(content)
                    chunks_with_pages = [{'text': chunk, 'page': None} for chunk in chunks]
                    logger.info(f"Successfully read {len(chunks_with_pages)} chunks from {file_path}")
                    return chunks_with_pages
            return []
        except Exception as e:
            logger.error(f"Error reading TXT {file_path}: {str(e)}")
            return []
    except Exception as e:
        logger.error(f"Error reading TXT {file_path}: {str(e)}")
        return []

# Improved document chunking function
def chunk_text(text, max_length=1000, min_length=300):
    if not text.strip():
        logger.warning("Empty text provided to chunk_text")
        return []
    
    # Split text into sentences
    sentences = sent_tokenize(text)
    chunks = []
    current_chunk = []
    current_length = 0
    
    for sentence in sentences:
        sentence_length = len(sentence) + 1
        if current_length + sentence_length > max_length and current_length >= min_length:
            chunks.append(" ".join(current_chunk))
            current_chunk = [sentence]
            current_length = sentence_length
        else:
            current_chunk.append(sentence)
            current_length += sentence_length
    
    if current_chunk and current_length >= min_length:
        chunks.append(" ".join(current_chunk))
    elif current_chunk:
        if chunks:
            chunks[-1] = chunks[-1] + " " + " ".join(current_chunk)
        else:
            chunks.append(" ".join(current_chunk))
    
    for i, chunk in enumerate(chunks):
        logger.debug(f"Chunk {i}: Length={len(chunk)}, Content={chunk[:100]}...")
    
    logger.info(f"Created {len(chunks)} chunks from text of length {len(text)}")
    return chunks

# Function to initialize the FAISS index
def initialize_index():
    global documents, embeddings, chunk_metadata, faiss_index, INDEX_INITIALIZED
    if INDEX_INITIALIZED:
        logger.info("Index already initialized, skipping...")
        return
    
    documents_dir = "documents"
    
    if not os.path.exists(documents_dir):
        os.makedirs(documents_dir)
        logger.warning(f"Created empty documents directory: {documents_dir}")
    
    documents = []
    embeddings = []
    chunk_metadata = []
    
    file_extractors = {
        '.txt': extract_text_from_txt,
        '.pdf': extract_text_from_pdf,
        '.pptx': extract_text_from_ppt,
        '.docx': extract_text_from_docx
    }
    
    logger.info("Starting document indexing...")
    for root, dirs, files in os.walk(documents_dir):
        for filename in files:
            ext = Path(filename).suffix.lower()
            if ext in file_extractors:
                filepath = os.path.join(root, filename)
                chunks_with_pages = file_extractors[ext](filepath)
                for i, chunk_info in enumerate(chunks_with_pages):
                    documents.append(chunk_info['text'])
                    chunk_metadata.append({
                        'filepath': filepath,
                        'chunk_index': i,
                        'original_text': chunk_info['text'],
                        'page': chunk_info['page']
                    })
    
    if documents:
        logger.info(f"Generating embeddings for {len(documents)} document chunks")
        embeddings = embedder.encode(documents, show_progress_bar=True)
        embeddings = np.array(embeddings).astype('float32')
    
        dimension = embeddings.shape[1]
        logger.info(f"Creating FAISS index with dimension {dimension}")
        faiss_index = faiss.IndexFlatL2(dimension)
        faiss_index.add(embeddings)
        logger.info(f"Indexed {len(documents)} document chunks")
        if len(documents) != len(chunk_metadata):
            logger.error(f"Mismatch between documents ({len(documents)}) and chunk_metadata ({len(chunk_metadata)})")
    else:
        logger.warning("No documents found to index. Please add files to the 'documents/' directory.")
        faiss_index = None
    
    INDEX_INITIALIZED = True

# Call the initialization function
initialize_index()

# Prompt template for user prompt
prompt_template = """Base your answer strictly on the information provided in the following internal project documents. Structure your answer as follows:
1. Start with a brief summary paragraph that provides an overview of the answer.
2. Follow with detailed information in a natural format. Use bullet points only when listing multiple items (e.g., information about multiple projects). Otherwise, use natural paragraphs.
3. For every piece of information you use from the documents, you MUST include the corresponding Reference ID (e.g., [Ref1], [Ref2]) inline with the text. If you do not use any document, explicitly state that no documents were used.
4. At the end of your answer, list the Reference IDs you used (e.g., References: [Ref1], [Ref2]).

Documents:
{passages}

Query: {query}

Answer:"""

def create_prompt(passages, query, max_context_tokens=3000):
    total_length = len(query) + len(prompt_template) - len("{passages}") - len("{query}")
    passages_text = ""
    passage_refs = []
    for idx, passage in enumerate(passages, 1):
        ref_id = f"[Ref{idx}]"
        if total_length + len(passage) < max_context_tokens * 4:
            passages_text += f"{ref_id} {passage}\n"
            passage_refs.append((ref_id, passage))
            total_length += len(passage) + 1
        else:
            logger.warning(f"Truncated passages to fit within {max_context_tokens} tokens")
            break
    logger.info(f"Passages provided in prompt: {passages_text}")
    return prompt_template.format(passages=passages_text, query=query), passage_refs

# Function to cluster passages based on semantic similarity
def cluster_passages(passages, embeddings, max_clusters=3):
    if len(passages) <= 1:
        return [(passages, embeddings)] if passages else []
    
    # Use KMeans to cluster embeddings
    num_clusters = min(max_clusters, len(passages))
    kmeans = KMeans(n_clusters=num_clusters, random_state=42)
    labels = kmeans.fit_predict(embeddings)
    
    # Group passages by cluster
    clustered_passages = [[] for _ in range(num_clusters)]
    clustered_embeddings = [[] for _ in range(num_clusters)]
    for idx, label in enumerate(labels):
        clustered_passages[label].append(passages[idx])
        clustered_embeddings[label].append(embeddings[idx])
    
    # Return clusters as list of (passages, embeddings) tuples
    return [(clustered_passages[i], clustered_embeddings[i]) for i in range(num_clusters) if clustered_passages[i]]

# Function to compute semantic similarity between two texts
def compute_similarity(text1, text2):
    embeddings = embedder.encode([text1, text2], show_progress_bar=False)
    similarity = np.dot(embeddings[0], embeddings[1]) / (np.linalg.norm(embeddings[0]) * np.linalg.norm(embeddings[1]))
    return similarity

# Function to format the response as HTML with citations at the end
def format_response(bot_message, passage_refs, merged_chunks):
    logger.info(f"Bot message: {bot_message}")
    
    # Convert the bot message to HTML, preserving paragraphs and bullet points
    lines = bot_message.strip().split('\n')
    formatted_lines = []
    in_list = False
    used_refs = set()
    
    # Parse the answer to find referenced passages
    for line in lines:
        line = line.strip()
        # Skip the "References:" line at the end
        if line.startswith("References:"):
            continue
        # Look for [RefX] patterns in the line
        refs_in_line = re.findall(r'\[Ref\d+\]', line)
        used_refs.update(refs_in_line)
        # Remove [RefX] from the line for display
        line = re.sub(r'\[Ref\d+\]', '', line).strip()
        if not line:
            if in_list:
                formatted_lines.append('</ul>')
                in_list = False
            continue
        if line.startswith('- '):
            if not in_list:
                formatted_lines.append('<ul>')
                in_list = True
            formatted_lines.append(f'<li>{line[2:]}</li>')
        else:
            if in_list:
                formatted_lines.append('</ul>')
                in_list = False
            formatted_lines.append(f'<p>{line}</p>')
    
    if in_list:
        formatted_lines.append('</ul>')
    
    # Generate citations for used references
    citations = []
    if used_refs:
        for ref_id, passage in passage_refs:
            if ref_id in used_refs:
                filepath = passage.split('\n')[0].replace("Document: ", "")
                for chunk in merged_chunks:
                    if chunk['filepath'] == filepath:
                        citation = f"Source: {os.path.basename(chunk['filepath'])}"
                        if chunk['pages']:
                            pages_str = ', '.join(map(str, sorted(set(chunk['pages']))))
                            citation += f", Pages/Slides: {pages_str}"
                        if citation not in citations:
                            citations.append(citation)
    else:
        # Fallback: Use semantic similarity to determine relevant passages
        logger.warning("No references explicitly used in the answer. Using semantic similarity to find relevant passages.")
        answer_text = bot_message.lower()
        for ref_id, passage in passage_refs:
            passage_text = passage.lower()
            # Check for keyword overlap
            answer_words = set(answer_text.split())
            passage_words = set(passage_text.split())
            common_words = answer_words.intersection(passage_words)
            # Compute semantic similarity
            similarity = compute_similarity(answer_text, passage_text)
            # Include the passage if there is significant overlap or high similarity
            if len(common_words) > 3 or similarity > 0.7:
                used_refs.add(ref_id)
                filepath = passage.split('\n')[0].replace("Document: ", "")
                for chunk in merged_chunks:
                    if chunk['filepath'] == filepath:
                        citation = f"Source: {os.path.basename(chunk['filepath'])}"
                        if chunk['pages']:
                            pages_str = ', '.join(map(str, sorted(set(chunk['pages']))))
                            citation += f", Pages/Slides: {pages_str}"
                        if citation not in citations:
                            citations.append(citation)
        if not citations:
            logger.warning("No relevant passages found via semantic similarity. No citations will be included.")

    # Add citations at the end
    if citations:
        citations_html = '<div class="references"><strong>References:</strong><br>' + '<br>'.join(citations) + '</div>'
        formatted_lines.append(citations_html)
    
    return ''.join(formatted_lines)


# Load configuration
try:
    with open('prompt.yaml', 'r', encoding='utf-8') as config_file:
        config = yaml.safe_load(config_file)
    SYSTEM_PROMPT = config.get('system_prompt', '')
    if not SYSTEM_PROMPT:
        raise ValueError("System prompt not found in prompt.yaml")
except FileNotFoundError:
    logger.error("prompt.yaml not found")
    raise
except yaml.YAMLError:
    logger.error("Invalid YAML in prompt.yaml")
    raise

@app.route('/')
def home():
    # Reset conversation history when the user visits the home page
    global conversation_history
    conversation_history = []
    return render_template('chat.html')

@app.route('/chat', methods=['POST'])
def chat():
    global conversation_history
    user_message = request.json['message']
    
    # Add user message to conversation history
    conversation_history.append({"role": "user", "content": user_message})
    
    try:
        query_embedding = embedder.encode([user_message], show_progress_bar=False)
        query_embedding = np.array(query_embedding).astype('float32')
        if len(query_embedding.shape) == 1:
            query_embedding = query_embedding.reshape(1, -1)
        logger.info(f"Query embedding shape: {query_embedding.shape}")
        
        passages = []
        passage_embeddings = []
        merged_chunks = []
        clusters = []
        if faiss_index is not None and len(documents) > 0:
            distances, indices = faiss_index.search(query_embedding, k=5)
            logger.info(f"Retrieved indices: {indices[0].tolist()}")
            logger.info(f"Distances: {distances[0].tolist()}")
            valid_indices = [idx for idx in indices[0] if 0 <= idx < len(chunk_metadata)]
            if not valid_indices:
                logger.warning("No valid indices retrieved from FAISS search")
                passages = []
                passage_embeddings = []
            else:
                valid_indices.sort(key=lambda idx: (chunk_metadata[idx]['filepath'], chunk_metadata[idx]['chunk_index']))
                
                current_chunk = None
                for idx in valid_indices:
                    chunk_info = chunk_metadata[idx]
                    if (current_chunk is None or
                            current_chunk['filepath'] != chunk_info['filepath'] or
                            current_chunk['chunk_index'] + 1 != chunk_info['chunk_index']):
                        if current_chunk is not None:
                            merged_chunks.append(current_chunk)
                        current_chunk = {
                            'filepath': chunk_info['filepath'],
                            'chunk_index': chunk_info['chunk_index'],
                            'text': chunk_info['original_text'],
                            'pages': [chunk_info['page']] if chunk_info['page'] is not None else []
                        }
                    else:
                        current_chunk['text'] += " " + chunk_info['original_text']
                        current_chunk['chunk_index'] = chunk_info['chunk_index']
                        if chunk_info['page'] is not None:
                            current_chunk['pages'].append(chunk_info['page'])
                if current_chunk is not None:
                    merged_chunks.append(current_chunk)
                
                # Create passages and embeddings for clustering
                for chunk in merged_chunks:
                    passage = f"Document: {chunk['filepath']}\nContent: {chunk['text']}"
                    passages.append(passage)
                    passage_embeddings.append(embedder.encode([passage], show_progress_bar=False)[0])
                
                # Cluster passages based on semantic similarity
                passage_embeddings = np.array(passage_embeddings).astype('float32')
                clusters = cluster_passages(passages, passage_embeddings)
                
                # Rebuild passages based on clusters
                passages = []
                for clustered_passages, _ in clusters:
                    cluster_text = "\n\n".join(clustered_passages)
                    passages.append(cluster_text)
                
                logger.info(f"Clustered passages: {passages}")
        else:
            passages = []
            passage_embeddings = []
            logger.warning("No documents indexed, proceeding without context")
    
    except Exception as e:
        logger.error(f"Embedding/search error: {str(e)}")
        passages = []
    
    # Generate prompt using the unified prompt template, and get passage references
    prompt, passage_refs = create_prompt(passages, user_message, max_context_tokens=3000)
    
    # Define the messages structure, including conversation history
    messages = [
        {
            "role": "system",
            "content": SYSTEM_PROMPT
        }
    ]
    # Add conversation history to messages
    messages.extend(conversation_history)
    # Add the current prompt as the latest user message
    messages.append({"role": "user", "content": prompt})
    
    # API call and response handling
    try:
        if API_TYPE == "deepseek":
            headers = {
                "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
                "Content-Type": "application/json"
            }
            payload = {
                "model": "deepseek-chat",
                "messages": messages,
                "max_tokens": 1500,
                "temperature": 1.3
            }
            deepseek_url = f"{DEEPSEEK_ENDPOINT}/chat/completions"
            logger.info(f"Making DeepSeek API call to: {deepseek_url}")
            response = requests.post(
                deepseek_url,
                headers=headers,
                json=payload,
                timeout=30
            )
            response.raise_for_status()
            
            response_data = response.json()
            if "choices" not in response_data or not response_data["choices"]:
                raise ValueError("No choices in DeepSeek response")
            bot_message = response_data["choices"][0].get("message", {}).get("content", "No response from DeepSeek").strip()

        elif API_TYPE == "azure":
            response = azure_client.chat.completions.create(
                model=AZURE_OPENAI_DEPLOYMENT,
                messages=messages,
                max_tokens=1500,
                temperature=0.4
            )
            bot_message = response.choices[0].message.content.strip()

        else:
            bot_message = "Error: API type not configured."

        # Format the response with citations at the end
        formatted_message = format_response(bot_message, passage_refs, merged_chunks)

        # Add bot response to conversation history
        conversation_history.append({"role": "assistant", "content": bot_message})

    except requests.exceptions.HTTPError as http_err:
        logger.error(f"HTTP error occurred: {http_err}")
        formatted_message = f"<p>Error: Failed to connect to the API. Please check the API endpoint ({DEEPSEEK_ENDPOINT}/chat/completions) and ensure it is correct.</p>"
    except Exception as e:
        logger.error(f"API call error: {str(e)}")
        formatted_message = f"<p>Error: {str(e)}</p>"
    
    return jsonify({"message": formatted_message})

if __name__ == '__main__':
    app.run(debug=True)